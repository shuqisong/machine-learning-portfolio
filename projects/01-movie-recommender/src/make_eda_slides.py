"""Generate a PowerPoint summarizing the Stage-1 EDA results.

Recomputes the figures/numbers from 01_exploration.ipynb and assembles them
into reports/eda_summary.pptx. Run with:
    .venv/bin/python projects/01-movie-recommender/src/make_eda_slides.py
"""
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt

PROJECT = Path(__file__).resolve().parent.parent
RAW = PROJECT / "data" / "raw"
FIGURES = PROJECT / "reports" / "figures"
FIGURES.mkdir(parents=True, exist_ok=True)

sns.set_theme(style="whitegrid")

movies = pd.read_csv(RAW / "movies.csv")
ratings = pd.read_csv(RAW / "ratings.csv")

n_users = ratings["userId"].nunique()
n_movies = ratings["movieId"].nunique()
n_ratings = len(ratings)
avg_rating = ratings["rating"].mean()

ratings_per_user = ratings.groupby("userId").size()
ratings_per_movie = ratings.groupby("movieId").size()
median_per_user = int(ratings_per_user.median())
median_per_movie = int(ratings_per_movie.median())

density = n_ratings / (n_users * n_movies)
sparsity = 1 - density

movie_stats = (
    ratings.groupby("movieId")
    .agg(num_ratings=("rating", "size"), avg_rating=("rating", "mean"))
    .merge(movies[["movieId", "title"]], on="movieId")
)
top_rated = movie_stats.sort_values("num_ratings", ascending=False).head(10)

# ---- Figure 1: rating distribution ----
fig, ax = plt.subplots(figsize=(7, 4.2))
sns.countplot(data=ratings, x="rating", color="steelblue", ax=ax)
ax.set_title("Distribution of movie ratings")
ax.set_xlabel("Rating (stars)")
ax.set_ylabel("Number of ratings")
fig.tight_layout()
fig1_path = FIGURES / "rating_distribution.png"
fig.savefig(fig1_path, dpi=150)
plt.close(fig)

# ---- Figure 2: ratings per user / per movie ----
fig, axes = plt.subplots(1, 2, figsize=(9.5, 4))
axes[0].hist(ratings_per_user, bins=50, color="steelblue")
axes[0].set_title("Ratings per user")
axes[0].set_xlabel("# ratings")
axes[0].set_ylabel("# users")
axes[1].hist(ratings_per_movie, bins=50, color="indianred")
axes[1].set_title("Ratings per movie")
axes[1].set_xlabel("# ratings")
axes[1].set_ylabel("# movies")
fig.tight_layout()
fig2_path = FIGURES / "activity_long_tail.png"
fig.savefig(fig2_path, dpi=150)
plt.close(fig)

# ---- Build the deck ----
prs = Presentation()
blank = prs.slide_layouts[6]
title_layout = prs.slide_layouts[0]

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
    return slide


def add_image_slide(title, image_path, caption=None):
    slide = prs.slides.add_slide(blank)
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), SLIDE_W - Inches(0.8), Inches(0.8))
    tf = tx.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    pic = slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.2), width=SLIDE_W - Inches(1.2))
    if caption:
        cap_top = Inches(1.2) + pic.height + Inches(0.15)
        cap = slide.shapes.add_textbox(Inches(0.6), cap_top, SLIDE_W - Inches(1.2), Inches(0.8))
        cap.text_frame.text = caption
        cap.text_frame.paragraphs[0].font.size = Pt(14)
        cap.text_frame.paragraphs[0].font.italic = True
    return slide


def add_table_slide(title, df, caption=None):
    slide = prs.slides.add_slide(blank)
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), SLIDE_W - Inches(0.8), Inches(0.8))
    tf = tx.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    rows, cols = df.shape[0] + 1, df.shape[1]
    table_w = SLIDE_W - Inches(1.2)
    table_h = Inches(0.4 * rows)
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.3), table_w, table_h)
    table = table_shape.table

    for c, col_name in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    for r in range(df.shape[0]):
        for c in range(df.shape[1]):
            cell = table.cell(r + 1, c)
            cell.text = str(df.iat[r, c])
            cell.text_frame.paragraphs[0].font.size = Pt(13)

    if caption:
        cap_top = Inches(1.3) + table_h + Inches(0.2)
        cap = slide.shapes.add_textbox(Inches(0.6), cap_top, SLIDE_W - Inches(1.2), Inches(0.8))
        cap.text_frame.text = caption
        cap.text_frame.paragraphs[0].font.size = Pt(14)
        cap.text_frame.paragraphs[0].font.italic = True
    return slide


# Slide 1: title
add_title_slide(
    "Movie Recommender — Stage 1: Exploratory Data Analysis",
    "MovieLens (small, latest) dataset  •  Generated from 01_exploration.ipynb",
)

# Slide 2: dataset overview
add_bullet_slide(
    "Dataset overview",
    [
        f"{n_ratings:,} ratings across {n_users:,} users and {n_movies:,} movies",
        f"Average rating: {avg_rating:.2f} / 5.0",
        f"Median ratings per user: {median_per_user}",
        f"Median ratings per movie: {median_per_movie}",
        "Source files: movies.csv (movieId, title, genres), ratings.csv (userId, movieId, rating, timestamp)",
    ],
)

# Slide 3: rating distribution
add_image_slide(
    "How are ratings distributed?",
    fig1_path,
    caption=(
        f"Ratings skew toward 3-4 stars (average = {avg_rating:.2f}). "
        "People tend to rate things they chose to watch, so high scores dominate."
    ),
)

# Slide 4: long tail
add_image_slide(
    "How active are users? How popular are movies?",
    fig2_path,
    caption=(
        f"Both distributions are long-tailed: median user rates {median_per_user} movies, "
        f"median movie gets {median_per_movie} ratings — a few power users/blockbusters dominate the rest."
    ),
)

# Slide 5: top rated movies table
top_table = top_rated[["title", "num_ratings", "avg_rating"]].copy()
top_table["avg_rating"] = top_table["avg_rating"].round(2)
top_table.columns = ["Title", "# Ratings", "Avg Rating"]
add_table_slide(
    "Most-rated movies (popularity baseline candidates)",
    top_table,
    caption="The most-rated titles are the natural seed for Stage 2's popularity-based recommender.",
)

# Slide 6: sparsity
add_bullet_slide(
    "How sparse is the data?",
    [
        f"User x Movie grid: {n_users:,} x {n_movies:,} = {n_users * n_movies:,} possible cells",
        f"Actual ratings collected: {n_ratings:,}",
        f"Density: {density:.2%}  (cells filled)",
        f"Sparsity: {sparsity:.2%}  (cells empty)",
        "This sparsity is why collaborative filtering needs to generalize from very few signals per user/movie.",
    ],
)

# Slide 7: takeaways / next steps
add_bullet_slide(
    "What we learned & what's next",
    [
        "Ratings skew toward 3-4 stars — people mostly rate things they like",
        "Users and movies both follow a long-tail distribution of activity/popularity",
        f"The data is extremely sparse ({sparsity:.1%}) — models must generalize from few signals",
        "Next: build a popularity baseline (02_baseline.ipynb), then content-based and collaborative filtering (SVD)",
    ],
)

out_path = PROJECT / "reports" / "eda_summary.pptx"
prs.save(out_path)
print(f"Saved deck to {out_path}")
print(f"Figures saved to {FIGURES}")
